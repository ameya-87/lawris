"""Generate Lawris_Full.pptx — 15-slide comprehensive pitch + Q&A deck.

Pairs with PITCH_DECK.md. Covers architecture, data model, all four pillars,
the three-iteration RAG evolution, structured citations, security, and the
competitive + roadmap slides.

Palette and helpers mirror build_pptx.py for visual consistency.
Run: python3 presentation/build_full_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Brand palette (shared with build_pptx.py) ────────────────────
INDIGO    = RGBColor(0x3F, 0x3D, 0x99)
INDIGO_D  = RGBColor(0x1E, 0x1B, 0x4B)
INDIGO_L  = RGBColor(0xC7, 0xC4, 0xF0)
STONE_50  = RGBColor(0xFA, 0xFA, 0xF9)
STONE_100 = RGBColor(0xF5, 0xF5, 0xF4)
STONE_200 = RGBColor(0xE7, 0xE5, 0xE4)
STONE_500 = RGBColor(0x78, 0x71, 0x6C)
STONE_900 = RGBColor(0x1C, 0x19, 0x17)
RED       = RGBColor(0xB9, 0x1C, 0x1C)
AMBER     = RGBColor(0xB4, 0x53, 0x09)
EMERALD   = RGBColor(0x05, 0x96, 0x69)
BLUE      = RGBColor(0x1D, 0x4E, 0xD8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# ─── Deck setup ───────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
TOTAL_SLIDES = 15

# ─── Helpers ──────────────────────────────────────────────────────
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    return bg

def add_text(slide, text, left, top, width, height,
             size=20, bold=False, color=STONE_900, align=PP_ALIGN.LEFT,
             font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf

def add_multiline(slide, lines, left, top, width, height,
                  size=14, color=STONE_900, bold=False,
                  font="Calibri", align=PP_ALIGN.LEFT, line_spacing=1.15):
    """Each entry in `lines` becomes a separate paragraph."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tf

def add_pill(slide, text, left, top, fill, fg=WHITE, size=11, width=1.6):
    h = Inches(0.32); w = Inches(width)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    pill.adjustments[0] = 0.5
    pill.line.fill.background()
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    tf = pill.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    return pill

def add_box(slide, left, top, width, height, fill=WHITE, border=STONE_200,
            border_width=0.75):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    b.fill.solid()
    b.fill.fore_color.rgb = fill
    b.line.color.rgb = border
    b.line.width = Pt(border_width)
    return b

def add_rounded_box(slide, left, top, width, height, fill=WHITE,
                    border=STONE_200, border_width=0.75, radius=0.06):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    b.adjustments[0] = radius
    b.fill.solid()
    b.fill.fore_color.rgb = fill
    b.line.color.rgb = border
    b.line.width = Pt(border_width)
    return b

def add_arrow(slide, x1, y1, x2, y2, color=STONE_500, width=1.5):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn

def add_footer(slide, n):
    add_text(slide, "LAWRIS — Full Deck", Inches(0.5), Inches(7.1),
             Inches(3), Inches(0.3), size=10, bold=True, color=STONE_500)
    add_text(slide, f"{n} / {TOTAL_SLIDES}", Inches(11.8), Inches(7.1),
             Inches(1), Inches(0.3), size=10, color=STONE_500, align=PP_ALIGN.RIGHT)

def add_section_label(slide, label):
    add_text(slide, label, Inches(0.75), Inches(0.55),
             Inches(10), Inches(0.5),
             size=13, bold=True, color=INDIGO)

def add_title(slide, text, top=1.1, size=36):
    add_text(slide, text, Inches(0.75), Inches(top),
             Inches(12), Inches(0.9),
             size=size, bold=True, color=STONE_900)

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, INDIGO_D)
add_text(s, "lawris.", Inches(0.7), Inches(1.9),
         Inches(10), Inches(2.2), size=128, bold=True, color=WHITE)
add_text(s, "An AI agent for Indian advocates.",
         Inches(0.75), Inches(3.95), Inches(11), Inches(0.8),
         size=28, color=INDIGO_L)
add_text(s,
         "Case management · deadline intelligence · AI drafting · case-grounded research",
         Inches(0.75), Inches(4.6), Inches(12), Inches(0.6),
         size=16, color=RGBColor(0x9C, 0x9A, 0xC4))
add_text(s, "Full pitch + Q&A deck", Inches(0.75), Inches(6.55),
         Inches(6), Inches(0.4),
         size=13, color=RGBColor(0x9C, 0x9A, 0xC4))
add_text(s, "Pairs with PITCH_DECK.md",
         Inches(7.5), Inches(6.55), Inches(5.5), Inches(0.4),
         size=13, color=RGBColor(0x9C, 0x9A, 0xC4), align=PP_ALIGN.RIGHT)
add_notes(s,
"COLD OPEN (memorised verbatim):\n"
"\"An accused person in India has 90 days. If their lawyer misses the "
"chargesheet deadline by even a day, they walk free. Lawyers in this country "
"carry 50 to 200 active cases each — on paper, WhatsApp, and three different "
"research portals. Calendaring failures drive one in four malpractice claims. "
"We built an AI agent that makes those misses impossible — and when the "
"lawyer asks what BNSS 187(3) says, it answers from the actual statute text, "
"cites the real Supreme Court ruling, and links straight to indiacode.nic.in.\"\n\n"
"Pause. Click.")

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 2)
add_section_label(s, "The problem")
add_title(s, "Indian lawyers are drowning in their own workflow.")

stats = [
    ("6+ hrs/week", "lost just locating documents",                 RED),
    ("40–60%",      "of billable time on repetitive drafts",         AMBER),
    ("77%",         "of lawyers run cases out of email + Word",      INDIGO),
    ("~25%",        "of malpractice claims trace to calendar misses", RED),
]
for i, (big, small, color) in enumerate(stats):
    x = Inches(0.75 + (i % 2) * 6.2)
    y = Inches(2.6 + (i // 2) * 1.75)
    add_text(s, big, x, y, Inches(5.8), Inches(1.0),
             size=46, bold=True, color=color)
    add_text(s, small, x, y + Inches(1.0), Inches(5.8), Inches(0.6),
             size=16, color=STONE_500)

add_text(s,
"No existing Indian tool unifies case management, AI drafting, deadline intelligence, "
"and research grounded in Indian law. MyKase has no AI. Harvey doesn't speak Indian law. "
"IndianKanoon is search-only. BNSS replaced CrPC in July 2024 — every template library is voided.",
         Inches(0.75), Inches(6.1), Inches(12), Inches(1.0),
         size=14, color=STONE_500)

add_notes(s,
"Read the four stats at pace — 5s each. Don't editorialise.\n"
"After stat four, the closing paragraph is the setup for 'we built it'.")

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution (4 pillars)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 3)
add_section_label(s, "What we built")
add_title(s, "Lawris — four pillars, one case file", size=34)
add_text(s,
"Open a case → auto-computed statutory deadlines → court-ready drafts in <30s → "
"research answers grounded in the case's own documents and Indian statutes.",
         Inches(0.75), Inches(2.0), Inches(12), Inches(1.0),
         size=18, color=STONE_900)

pillars = [
    ("DEADLINE BRAIN",
     "BNSS 187(3) chargesheet rule, Limitation Act periods, hearing alerts — all auto-generated from FIR data.",
     RED),
    ("AI DRAFTING",
     "Bail applications (§483 BNSS) and plaints (Order VII CPC) streamed in under 30 seconds with real Indian-court formatting.",
     INDIGO),
    ("HEARING MEMORY",
     "Every logged hearing regenerates the case summary. Next-date hearings become deadline rows automatically.",
     EMERALD),
    ("GROUNDED RESEARCH",
     "Answers cite the lawyer's own uploads AND a curated corpus of Indian statutes + landmark judgments.",
     BLUE),
]
for i, (t, body, color) in enumerate(pillars):
    x = Inches(0.75 + (i % 2) * 6.2)
    y = Inches(3.4 + (i // 2) * 1.7)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.15), Inches(0.4))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = color
    add_text(s, t, x + Inches(0.3), y, Inches(5.6), Inches(0.4),
             size=13, bold=True, color=color)
    add_text(s, body, x + Inches(0.3), y + Inches(0.5), Inches(5.6), Inches(1.0),
             size=13, color=STONE_900)

add_notes(s,
"Four pillars — the first three are the Round 1 product; the fourth (grounded "
"research) is where most of the Round 2 innovation lives. Drafting is the demo "
"spectacle; research is the technical differentiator.")

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 4)
add_section_label(s, "System architecture")
add_title(s, "One Next.js app · Supabase · Gemini", size=28)

# Browser
add_rounded_box(s, Inches(1.0), Inches(1.9), Inches(11.3), Inches(0.9),
                fill=INDIGO_D, border=INDIGO_D)
add_text(s, "BROWSER  ·  React Server Components + Client Components + AI stream reader",
         Inches(1.2), Inches(2.15), Inches(11), Inches(0.5),
         size=14, bold=True, color=WHITE)

# Next.js box
add_rounded_box(s, Inches(1.0), Inches(3.05), Inches(11.3), Inches(2.1),
                fill=WHITE, border=INDIGO)
add_text(s, "NEXT.JS 14 — App Router + API handlers + Middleware",
         Inches(1.2), Inches(3.15), Inches(11), Inches(0.35),
         size=13, bold=True, color=INDIGO)

routes = [
    ("Pages", "/, /cases, /cases/:id, /calendar"),
    ("/api/ai/draft", "streaming Gemini 2.5 Flash"),
    ("/api/ai/research", "hybrid RAG + structured JSON + URL enrichment"),
    ("/api/documents/ingest", "pdf-parse → chunk → embed → Storage"),
    ("/api/cases, /api/deadlines/*", "CRUD + auto-deadlines"),
    ("Middleware", "Supabase SSR session + LAWYER_ID fallback"),
]
for i, (left_text, right_text) in enumerate(routes):
    y = Inches(3.55 + (i // 2) * 0.5)
    x = Inches(1.3 + (i % 2) * 5.3)
    add_text(s, left_text, x, y, Inches(2.0), Inches(0.3),
             size=11, bold=True, color=STONE_900)
    add_text(s, right_text, x + Inches(2.1), y, Inches(3.2), Inches(0.3),
             size=11, color=STONE_500)

# Backends row
backend_y = Inches(5.45)
for i, (label, body, fill) in enumerate([
    ("SUPABASE",     "Postgres + pgvector · Storage · Auth", EMERALD),
    ("GOOGLE GEMINI", "2.5 Flash drafting · Flash-Lite research · embedding-001", AMBER),
    ("vendor helpers", "zod · date-fns · pdf-parse · tsx · dotenv", STONE_500),
]):
    x = Inches(1.0 + i * 3.85)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, backend_y,
                             Inches(0.15), Inches(1.0))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = fill
    add_text(s, label, x + Inches(0.25), backend_y, Inches(3.5), Inches(0.4),
             size=12, bold=True, color=fill)
    add_text(s, body, x + Inches(0.25), backend_y + Inches(0.4),
             Inches(3.5), Inches(0.7),
             size=11, color=STONE_900)

add_text(s,
"Architectural bets: monolithic Next.js (no microservices), pgvector "
"instead of Pinecone, streaming via ReadableStream (no WebSockets), "
"RSC-first data fetching with client components only where interactive.",
         Inches(1.0), Inches(6.75), Inches(11.3), Inches(0.6),
         size=11, color=STONE_500)

add_notes(s,
"Architecture slide is for engineer-judges. Non-technical judges skim it.\n"
"Headline: 'One deployable unit. No microservices. pgvector keeps embeddings "
"in the same database as the cases — no Pinecone, no extra SLA.'")

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Data Model
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 5)
add_section_label(s, "Data model")
add_title(s, "9 tables · 9 enums · pgvector embedded", size=28)

tables_left = [
    ("users",           "1 row",  ["id, email, full_name, bar_council_no"]),
    ("clients",         "2 rows", ["id, lawyer_id, full_name, aadhar_no"]),
    ("cases",           "3 rows", ["fir_date, sections, offence_max_years", "phase, status, ai_summary"]),
    ("deadlines",       "4 rows", ["case_id, due_date, urgency", "is_auto_generated"]),
    ("documents",       "varies", ["case_id, name, doc_type, storage_path"]),
]
tables_right = [
    ("hearing_logs",    "2 rows", ["hearing_date, what_happened", "judge_order, next_date"]),
    ("research_notes",  "grows",  ["query, content (JSON), citation, tags"]),
    ("document_chunks", "per case", ["case_id FK · vector(768) embedding", "used by retrieveChunks()"]),
    ("corpus_chunks",   "~761 rows", ["source_type (statute/judgment)", "source_name, category, vector(768)"]),
]
def draw_table(lst, x_base):
    for i, (name, count, fields) in enumerate(lst):
        y = Inches(1.95 + i * 0.95)
        add_rounded_box(s, x_base, y, Inches(5.9), Inches(0.85), fill=WHITE)
        add_text(s, name, x_base + Inches(0.2), y + Inches(0.1),
                 Inches(2.0), Inches(0.35),
                 size=13, bold=True, color=INDIGO, font="Consolas")
        add_text(s, count, x_base + Inches(2.2), y + Inches(0.1),
                 Inches(1.4), Inches(0.35),
                 size=11, color=STONE_500)
        for j, f in enumerate(fields):
            add_text(s, f, x_base + Inches(0.2), y + Inches(0.42 + j * 0.20),
                     Inches(5.6), Inches(0.22),
                     size=10, color=STONE_900, font="Consolas")

draw_table(tables_left,  Inches(0.65))
draw_table(tables_right, Inches(6.75))

add_text(s,
"All FKs cascade-delete from case or user. Enums: case_type, case_phase, case_status, "
"court_type, deadline_type, urgency, doc_type, doc_source, research_source.",
         Inches(0.65), Inches(6.85), Inches(12), Inches(0.5),
         size=11, color=STONE_500)

add_notes(s,
"9 tables. The two chunk tables (document_chunks for per-case, corpus_chunks for "
"shared Indian law) share the vector(768) embedding space — a 'law chunk' and a "
"'case chunk' are directly comparable. That's what makes the two-tier retrieval "
"cheap.")

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Pillar 1: Deadline Brain
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 6)
add_section_label(s, "Pillar 1 — Deadline Brain")
add_title(s, "Statutory law as code. Zero AI. Zero hallucination.",
          size=30)

# Left: the input
add_rounded_box(s, Inches(0.75), Inches(2.1), Inches(4.0), Inches(3.5),
                fill=WHITE)
add_text(s, "INPUT", Inches(0.95), Inches(2.25), Inches(3.5), Inches(0.3),
         size=11, bold=True, color=INDIGO)
add_multiline(s, [
    "POST /api/cases",
    "",
    "  title: 'Sharma — POCSO'",
    "  case_type: 'criminal'",
    "  fir_date: '2026-02-01'",
    "  offence_max_years: 10",
    "  sections: '506, 376 BNS'",
], Inches(0.95), Inches(2.65), Inches(3.7), Inches(2.8),
   size=12, color=STONE_900, font="Consolas")

# Middle: the rules
add_rounded_box(s, Inches(5.0), Inches(2.1), Inches(3.5), Inches(3.5),
                fill=INDIGO_D, border=INDIGO_D)
add_text(s, "RULES (deadlines.ts)", Inches(5.2), Inches(2.25),
         Inches(3.1), Inches(0.3),
         size=11, bold=True, color=WHITE)
add_multiline(s, [
    "BNSS §187(3):",
    "  <10y → 60d chargesheet",
    "  ≥10y → 90d chargesheet",
    "",
    "Limitation Act:",
    "  civil → 3y default",
    "",
    "Hearing logs:",
    "  next_date → new row",
], Inches(5.2), Inches(2.65), Inches(3.2), Inches(2.8),
   size=12, color=WHITE, font="Consolas")

# Right: the output
add_rounded_box(s, Inches(8.75), Inches(2.1), Inches(4.0), Inches(3.5),
                fill=WHITE)
add_text(s, "OUTPUT", Inches(8.95), Inches(2.25), Inches(3.5), Inches(0.3),
         size=11, bold=True, color=EMERALD)
add_multiline(s, [
    "deadlines rows:",
    "",
    "  2026-04-02  chargesheet  RED",
    "  2026-05-15  hearing      AMBER",
    "  2029-02-01  limitation   GREEN",
    "",
    "classifyUrgency(due_date)",
    "  critical ≤3d / high ≤7d",
    "  medium ≤30d / low >30d",
], Inches(8.95), Inches(2.65), Inches(3.7), Inches(2.8),
   size=11, color=STONE_900, font="Consolas")

# Arrows
add_arrow(s, Inches(4.75), Inches(3.85), Inches(5.0), Inches(3.85),
          color=INDIGO, width=2)
add_arrow(s, Inches(8.5), Inches(3.85), Inches(8.75), Inches(3.85),
          color=INDIGO, width=2)

add_text(s,
"UI surfacing: color-coded pills on dashboard, case detail, and calendar. "
"Red pill on the POCSO case is the thing the judges will see in the demo "
"— that's the deadline the system computed and the lawyer would otherwise have missed.",
         Inches(0.75), Inches(6.0), Inches(12), Inches(1.1),
         size=12, color=STONE_500)

add_notes(s,
"Say: 'Deadlines aren't AI. They're statute translated into TypeScript. "
"Zero hallucination risk. This is the feature judges should be MOST "
"confident about, not least.'")

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Pillar 2: Streaming AI Drafting
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 7)
add_section_label(s, "Pillar 2 — AI document drafting (streaming)")
add_title(s, "Tokens → ReadableStream → editor pane", size=28)

# Flow boxes
boxes = [
    (0.75,  "USER",      "Click\nDraft Bail", INDIGO),
    (3.0,   "/api/ai/draft", "zod parse\nload case+client", STONE_500),
    (5.5,   "Gemini 2.5 Flash", "generateContentStream\nprompt: bail-application.ts", AMBER),
    (8.5,   "ReadableStream", "for-await chunks\nenqueue bytes", INDIGO),
    (11.0,  "<AiDraftPanel>", "fetch reader\nsetDraft() each tick", EMERALD),
]
y_box = Inches(2.2)
for x_in, title, body, color in boxes:
    x = Inches(x_in)
    add_rounded_box(s, x, y_box, Inches(2.0), Inches(1.6), fill=WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_box,
                             Inches(2.0), Inches(0.1))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = color
    add_text(s, title, x + Inches(0.15), y_box + Inches(0.2),
             Inches(1.8), Inches(0.35),
             size=11, bold=True, color=color)
    add_multiline(s, body.split("\n"),
                  x + Inches(0.15), y_box + Inches(0.6),
                  Inches(1.75), Inches(0.95),
                  size=10, color=STONE_900, font="Consolas")
# Arrows between
for i, x_in in enumerate([2.75, 5.25, 8.25, 10.75]):
    add_arrow(s, Inches(x_in), Inches(3.0),
              Inches(x_in + 0.25), Inches(3.0),
              color=STONE_500, width=1.5)

# Bottom — what gets persisted
add_rounded_box(s, Inches(0.75), Inches(4.3), Inches(12.2), Inches(1.5),
                fill=WHITE)
add_text(s, "ON STREAM COMPLETE",
         Inches(0.95), Inches(4.45), Inches(6), Inches(0.3),
         size=11, bold=True, color=INDIGO)
add_multiline(s, [
    "INSERT documents (source: ai_drafted, content: full draft Markdown, ai_prompt_used: snapshot)",
    "",
    "Quality signal from live test: bail app cited Satender Kumar Antil + Dataram Singh + Balchand,",
    "flagged the chargesheet as overdue based on FIR date — WITHOUT being told explicitly.",
], Inches(0.95), Inches(4.8), Inches(11.8), Inches(0.9),
   size=11, color=STONE_900)

add_text(s,
"Models: Drafting = gemini-2.5-flash, Research/Summarise = gemini-2.5-flash-lite. "
"All three IDs configurable via env vars — swap to Pro without touching code.",
         Inches(0.75), Inches(6.4), Inches(12), Inches(0.7),
         size=11, color=STONE_500)

add_notes(s,
"In the demo: click Draft Bail, then SILENCE for 20 seconds while the "
"document writes itself. Don't narrate. Let the audience watch. The payoff "
"comes AFTER the stream finishes: 'Court-ready. Real precedents. 25 seconds.'")

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Pillar 3a: Per-case RAG (RAG v1)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 8)
add_section_label(s, "Research — iteration 1 · per-case RAG")
add_title(s, "Lawyer's own uploads → embedded → retrievable", size=28)

# Ingestion flow
y = Inches(2.1)
steps = [
    ("PDF upload", "multipart\nformData"),
    ("documents row", "UUID\nissued"),
    ("Supabase\nStorage", "case-documents/\n{case_id}/\n{doc.id}_..."),
    ("pdf-parse", "text layer\nextraction"),
    ("chunkText", "500 words\n50 overlap"),
    ("embedText", "gemini-\nembedding-001\n@ 768 dims"),
    ("document_chunks", "vector(768)\ncase_id scoped"),
]
step_w = 1.7
gap = 0.05
for i, (title, body) in enumerate(steps):
    x = Inches(0.6 + i * (step_w + gap))
    add_rounded_box(s, x, y, Inches(step_w), Inches(1.6),
                    fill=WHITE, border=INDIGO_L)
    add_text(s, title, x + Inches(0.1), y + Inches(0.15),
             Inches(step_w - 0.2), Inches(0.4),
             size=10, bold=True, color=INDIGO,
             align=PP_ALIGN.CENTER)
    add_multiline(s, body.split("\n"),
                  x + Inches(0.1), y + Inches(0.6),
                  Inches(step_w - 0.2), Inches(1.0),
                  size=9, color=STONE_900, align=PP_ALIGN.CENTER,
                  font="Consolas")
    if i < len(steps) - 1:
        add_arrow(s, x + Inches(step_w), y + Inches(0.8),
                  x + Inches(step_w + gap), y + Inches(0.8),
                  color=STONE_500, width=1.2)

# Retrieval
add_rounded_box(s, Inches(0.75), Inches(4.1), Inches(12.2), Inches(1.8),
                fill=INDIGO_D, border=INDIGO_D)
add_text(s, "RETRIEVAL",
         Inches(0.95), Inches(4.25), Inches(6), Inches(0.3),
         size=11, bold=True, color=INDIGO_L)
add_multiline(s, [
    "retrieveChunks(query, case_id, topK=3) →",
    "  embedText(query)                             // 768-dim query vector",
    "  RPC match_chunks(query_embedding, case_id)    // cosine similarity",
    "  returns top-3 ChunkMatch[] scoped to this case",
], Inches(0.95), Inches(4.6), Inches(11.9), Inches(1.2),
   size=12, color=WHITE, font="Consolas")

add_text(s,
"Signed URLs: src/lib/storage.ts creates a 1-hour signed URL for each uploaded PDF. "
"The research response includes case_documents[] so the UI can turn [DOC] citations "
"into clickable links that open the original PDF.",
         Inches(0.75), Inches(6.1), Inches(12), Inches(0.9),
         size=11, color=STONE_500)

add_notes(s,
"Seven-step ingestion pipeline. Key detail: the chunks table stores case_id as FK, "
"so retrieval is automatically scoped — one lawyer's uploads can't surface in "
"another's research.")

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Pillar 3b: Two-tier corpus (RAG v2)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 9)
add_section_label(s, "Research — iteration 2 · two-tier corpus")
add_title(s, "Per-case docs + shared Indian law corpus", size=28)

# Two columns: [per-case] and [corpus]
col_w = 5.8
for i, (title, rows, color) in enumerate([
    ("PER-CASE TIER (document_chunks)", [
        "Scope: one lawyer's case_id",
        "Source: their uploaded PDFs",
        "Ingestion: /api/documents/ingest",
        "Retrieval: retrieveChunks(q, case_id, 3)",
        "Prompt prefix: [DOC]",
        "UI badge: amber",
    ], AMBER),
    ("CORPUS TIER (corpus_chunks)", [
        "Scope: shared across all lawyers",
        "Source: 6 statutes + 8 landmark judgments",
        "Ingestion: scripts/ingest-corpus.ts",
        "Retrieval: retrieveCorpusChunksHybrid(q, 6)",
        "Prompt prefix: [LAW: source_name]",
        "UI badge: blue",
    ], BLUE),
]):
    x = Inches(0.75 + i * (col_w + 0.3))
    add_rounded_box(s, x, Inches(1.95), Inches(col_w), Inches(3.3),
                    fill=WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.95),
                             Inches(0.15), Inches(3.3))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = color
    add_text(s, title, x + Inches(0.3), Inches(2.1),
             Inches(col_w - 0.3), Inches(0.4),
             size=13, bold=True, color=color)
    add_multiline(s, rows, x + Inches(0.3), Inches(2.55),
                  Inches(col_w - 0.3), Inches(2.6),
                  size=12, color=STONE_900)

# Corpus contents block
add_rounded_box(s, Inches(0.75), Inches(5.4), Inches(12.2), Inches(1.7),
                fill=INDIGO_D, border=INDIGO_D)
add_text(s, "CORPUS CONTENTS · 14 sources · ~761 chunks",
         Inches(0.95), Inches(5.55), Inches(12), Inches(0.3),
         size=12, bold=True, color=INDIGO_L)
add_multiline(s, [
    "Statutes (6):  BNSS 2023  ·  BNS 2023  ·  BSA 2023  ·  CPC 1908  ·  POCSO 2012  ·  Human Rights Act 1993",
    "Judgments (8): Arnesh Kumar · Satender Kumar Antil · Maneka Gandhi · Hans Kumar ·",
    "               A.S. Templeton · Mahender Bansal · Rahul Pathak · Sushma Trivedi",
], Inches(0.95), Inches(5.9), Inches(11.9), Inches(1.1),
   size=11, color=WHITE, font="Consolas")

add_notes(s,
"Same vector(768) embedding space across both tables — chunks from a lawyer's FIR "
"and chunks from BNSS are directly comparable. The research route runs both "
"retrievals in Promise.all and lets Gemini see both blocks in the prompt.")

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — Pillar 3c: Hybrid retrieval (RAG v3)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 10)
add_section_label(s, "Research — iteration 3 · hybrid retrieval")
add_title(s, "The fix for section-number queries", size=28)

# Problem callout
add_rounded_box(s, Inches(0.75), Inches(1.95), Inches(12.2), Inches(1.1),
                fill=RED, border=RED)
add_text(s, "PROBLEM",
         Inches(0.95), Inches(2.1), Inches(6), Inches(0.3),
         size=11, bold=True, color=WHITE)
add_text(s,
"Query \"BNSS §187(3) chargesheet deadline\" → top hit was Satender Kumar Antil "
"(a JUDGMENT about §187) at sim 0.58. Actual statute text didn't crack top-3.",
         Inches(0.95), Inches(2.4), Inches(11.9), Inches(0.7),
         size=12, color=WHITE)

# Three-layer fix
layers = [
    ("LAYER 1 · keyword boost",
     "Regex extract section numbers from query\n→ ILIKE match statute chunks only\n→ similarity = 1.00 (artificial)",
     EMERALD),
    ("LAYER 2 · stratified semantic",
     "match_corpus_stratified RPC\n→ top-3 statutes + top-3 judgments\n→ prevents judgments monopolising",
     BLUE),
    ("LAYER 3 · merge + dedup",
     "Keyword hits ranked first\n→ semantic hits appended\n→ dedup by first-100-char signature\n→ cap at topK=6",
     INDIGO),
]
for i, (title, body, color) in enumerate(layers):
    x = Inches(0.75 + i * 4.1)
    y = Inches(3.25)
    add_rounded_box(s, x, y, Inches(3.95), Inches(2.0), fill=WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                             Inches(3.95), Inches(0.1))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = color
    add_text(s, title, x + Inches(0.2), y + Inches(0.2),
             Inches(3.7), Inches(0.4),
             size=12, bold=True, color=color)
    add_multiline(s, body.split("\n"),
                  x + Inches(0.2), y + Inches(0.65),
                  Inches(3.7), Inches(1.35),
                  size=11, color=STONE_900)

# Measured impact
add_rounded_box(s, Inches(0.75), Inches(5.45), Inches(12.2), Inches(1.6),
                fill=EMERALD, border=EMERALD)
add_text(s, "MEASURED IMPACT  ·  same query, before vs after",
         Inches(0.95), Inches(5.6), Inches(12), Inches(0.3),
         size=11, bold=True, color=WHITE)
add_multiline(s, [
    "BEFORE  corpus_sources: [judgment:Satender Kumar Antil, …]  ·  similarities: [0.58, 0.55, 0.54]",
    "AFTER   corpus_sources: [statute:Bharatiya Nyaya Sanhita,   ·  similarities: [1.00, 1.00, 1.00]",
    "                         statute:Bharatiya Nyaya Sanhita,",
    "                         statute:Bharatiya Nagarik Suraksha Sanhita]",
], Inches(0.95), Inches(5.95), Inches(11.9), Inches(1.1),
   size=11, color=WHITE, font="Consolas")

add_notes(s,
"This is the defensible technical differentiator. Plain semantic retrieval "
"ranks JUDGMENTS above STATUTES for section-number queries because judgment "
"prose literally says 'Section 187 chargesheet' while the statute uses "
"'sixty days' without the section number in the chunk. We caught it, fixed "
"it, measured it, and the fix is in production.")

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — Structured citations + source verification
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 11)
add_section_label(s, "Structured citations · source verification")
add_title(s, "Every citation is a card with a clickable link",
          size=28)

# Citation schema box
add_rounded_box(s, Inches(0.75), Inches(1.95), Inches(6.0), Inches(4.2),
                fill=WHITE)
add_text(s, "GEMINI RESPONSE SCHEMA (per citation)",
         Inches(0.95), Inches(2.1), Inches(5.6), Inches(0.3),
         size=11, bold=True, color=INDIGO)
add_multiline(s, [
    "{",
    "  source:             '[LAW: Satender Kumar Antil …]',",
    "  case_name_or_statute: 'Satender Kumar Antil v. CBI (2022)',",
    "  source_type:        'law' | 'doc',",
    "  core_holding:       '1-2 sentences…',",
    "  key_facts:          '1-2 sentences…',",
    "  relevance_to_query: '1-2 sentences…',",
    "",
    "  // Server-enriched after Gemini call:",
    "  source_url:         'https://indiankanoon.org/doc/…'",
    "}",
], Inches(0.95), Inches(2.5), Inches(5.8), Inches(3.5),
   size=11, color=STONE_900, font="Consolas")

# Rendered card example (right side)
card_x = Inches(7.0)
add_rounded_box(s, card_x, Inches(1.95), Inches(5.9), Inches(4.2),
                fill=WHITE)
add_text(s, "RENDERED UI",
         card_x + Inches(0.2), Inches(2.1), Inches(5.6), Inches(0.3),
         size=11, bold=True, color=INDIGO)
# Badge + clickable title row
add_pill(s, "LAW", card_x + Inches(0.25), Inches(2.5), BLUE,
         size=10, width=0.7)
add_text(s, "Satender Kumar Antil v. CBI (2022)  ↗",
         card_x + Inches(1.05), Inches(2.53), Inches(4.7), Inches(0.3),
         size=12, bold=True, color=RGBColor(0x4F, 0x46, 0xE5))
# Three labeled fields
field_y_start = 3.0
for i, (label, body) in enumerate([
    ("CORE HOLDING",
     "Default bail is indefeasible. Delay beyond statutory period compels release."),
    ("KEY FACTS",
     "Accused detained beyond 60 days; SC issued guidelines against routine extensions."),
    ("WHY THIS MATTERS HERE",
     "POCSO case crosses 60-day mark in 2 days — default-bail ground is live."),
]):
    y = Inches(field_y_start + i * 1.0)
    add_text(s, label, card_x + Inches(0.25), y,
             Inches(5.5), Inches(0.25),
             size=9, bold=True, color=STONE_500)
    add_text(s, body, card_x + Inches(0.25), y + Inches(0.25),
             Inches(5.5), Inches(0.7),
             size=11, color=STONE_900)

# Bottom: URL mapping + fuzzy match
add_rounded_box(s, Inches(0.75), Inches(6.3), Inches(12.2), Inches(0.85),
                fill=INDIGO_D, border=INDIGO_D)
add_text(s,
"URL enrichment: src/lib/source-urls.ts maps 14 corpus sources to indiacode.nic.in / indiankanoon.org. "
"Fuzzy token-based matching handles Gemini's name variations.  "
"DOC citations: fuzzy-match against case_documents[] from Supabase Storage signed URLs.",
         Inches(0.95), Inches(6.45), Inches(11.9), Inches(0.7),
         size=11, color=WHITE)

add_notes(s,
"Click a LAW citation in the demo — it opens indiankanoon.org or indiacode.nic.in "
"in a new tab. That's the 'every answer is verifiable' moment. "
"Fuzzy URL matching means small name variations (e.g. missing year) still resolve.")

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — Tech stack grid
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 12)
add_section_label(s, "Tech stack")
add_title(s, "Every dep · why chosen", size=30)

stack_groups = [
    ("Frontend", [
        ("Next.js 14.2",         "App Router, RSC, streaming"),
        ("React 18",             "UI runtime"),
        ("Tailwind 3.4",         "Utility CSS"),
        ("TanStack Query 5.99",  "Server-state cache"),
        ("Lucide",               "Icons"),
    ], INDIGO),
    ("Backend & data", [
        ("Supabase JS 2.103",    "Postgres + Storage client"),
        ("@supabase/ssr",        "Middleware session"),
        ("zod 4.3",              "Runtime validation"),
        ("date-fns 4.1",         "Deadline math"),
        ("pdf-parse v2",         "PDF text extraction"),
    ], EMERALD),
    ("AI", [
        ("@google/genai 1.50",   "Gemini SDK"),
        ("gemini-2.5-flash",     "Drafting"),
        ("gemini-2.5-flash-lite","Research + summarise"),
        ("gemini-embedding-001", "768-dim vectors"),
        ("ReadableStream",       "Token streaming"),
    ], AMBER),
    ("Scripts & tooling", [
        ("tsx 4.21",             "TS scripts (ingestion, probes)"),
        ("dotenv",               ".env.local loading"),
        ("python-pptx 1.0",      "This deck"),
        ("eslint (next config)", "Lint on build"),
        ("nohup / pkill",        "Dev-server lifecycle"),
    ], STONE_500),
]
col_w = 2.95
for i, (title, rows, color) in enumerate(stack_groups):
    x = Inches(0.65 + i * (col_w + 0.12))
    y = Inches(1.9)
    add_rounded_box(s, x, y, Inches(col_w), Inches(4.8), fill=WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                             Inches(col_w), Inches(0.1))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = color
    add_text(s, title, x + Inches(0.2), y + Inches(0.2),
             Inches(col_w - 0.3), Inches(0.4),
             size=13, bold=True, color=color)
    for j, (dep, purpose) in enumerate(rows):
        ry = y + Inches(0.8 + j * 0.75)
        add_text(s, dep, x + Inches(0.2), ry,
                 Inches(col_w - 0.4), Inches(0.3),
                 size=11, bold=True, color=STONE_900, font="Consolas")
        add_text(s, purpose, x + Inches(0.2), ry + Inches(0.3),
                 Inches(col_w - 0.4), Inches(0.4),
                 size=10, color=STONE_500)

add_text(s,
"All model IDs configurable via env — swap to paid/Pro tier without code changes. "
"Prompt files are vendor-neutral; the Gemini wrapper in src/lib/gemini.ts is the only "
"file that needs editing to switch LLM providers.",
         Inches(0.65), Inches(6.85), Inches(12), Inches(0.5),
         size=11, color=STONE_500)

add_notes(s,
"Engineer judges will ask 'why Gemini not GPT-4'. Answer: free tier + structured "
"output + long context. The SDK is isolated in one file — switch to OpenAI or "
"Anthropic in an afternoon.")

# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — Security & auth
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 13)
add_section_label(s, "Security & auth")
add_title(s, "Today's posture · production hardening path",
          size=28)

# Two columns
col_w = 5.9
cols = [
    ("SHIPPED TODAY", [
        ("Supabase Auth", "email+password, middleware validates session"),
        ("Service-role key", "server-only via src/lib/supabase/server.ts"),
        ("Zod validation", "every POST body parsed at the boundary"),
        ("Storage bucket private", "case-documents not public, signed URLs 1h TTL"),
        ("LAWYER_ID env pin", "demo-day fallback when session cookie dies"),
        ("Prompt grounding", "corpus-only citation rule + responseSchema"),
    ], EMERALD),
    ("POST-HACKATHON GAPS", [
        ("Supabase RLS", "every table gated by auth.uid() on lawyer_id"),
        ("Aadhar encryption", "pgsodium column encryption"),
        ("Audit log", "every /api/ai/* call logged with lawyer_id"),
        ("Rate limiting", "per-lawyer upstream Gemini quota enforcement"),
        ("Prompt injection tests", "adversarial case-notes suite"),
        ("SSO / bar council lookup", "verify real advocate before sign-up"),
    ], AMBER),
]
for i, (title, rows, color) in enumerate(cols):
    x = Inches(0.75 + i * (col_w + 0.3))
    add_rounded_box(s, x, Inches(1.95), Inches(col_w), Inches(5.0),
                    fill=WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.95),
                             Inches(0.15), Inches(5.0))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = color
    add_text(s, title, x + Inches(0.3), Inches(2.1),
             Inches(col_w - 0.3), Inches(0.4),
             size=13, bold=True, color=color)
    for j, (label, body) in enumerate(rows):
        y = Inches(2.6 + j * 0.73)
        add_text(s, label, x + Inches(0.3), y,
                 Inches(col_w - 0.5), Inches(0.3),
                 size=12, bold=True, color=STONE_900)
        add_text(s, body, x + Inches(0.3), y + Inches(0.3),
                 Inches(col_w - 0.5), Inches(0.4),
                 size=11, color=STONE_500)

add_notes(s,
"Answer to 'where's the security?' — we ship real Supabase Auth today. RLS "
"is the ONE thing we know we'll do day one after the demo. Everything else "
"is on a published roadmap, not a surprise.")

# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — Competitive landscape
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 14)
add_section_label(s, "Competitive positioning")
add_title(s, "Adjacent products solve part of this. None solve all.",
          size=26)

cols = ["", "Case mgmt", "AI drafting", "Indian law", "Deadline brain", "Hybrid RAG"]
rows = [
    ("Clio (US)",        "Yes",  "Basic", "—",   "Yes", "—"),
    ("MyKase (India)",   "Yes",  "—",     "Yes", "Yes", "—"),
    ("Harvey AI (US)",   "—",    "Yes",   "—",   "—",   "Keyword only"),
    ("IndianKanoon",     "—",    "—",     "Yes", "—",   "— (search only)"),
    ("Lawris",           "Yes",  "Yes",   "Yes", "Yes", "Yes"),
]
table_left = Inches(0.75)
table_top  = Inches(2.3)
col_widths = [Inches(2.5), Inches(1.95), Inches(1.95), Inches(1.95),
              Inches(2.1), Inches(1.95)]
row_h = Inches(0.6)
table = s.shapes.add_table(len(rows) + 1, len(cols),
                           table_left, table_top,
                           sum(col_widths, Emu(0)),
                           row_h * (len(rows) + 1)).table
for i, w in enumerate(col_widths):
    table.columns[i].width = w
for r in range(len(rows) + 1):
    table.rows[r].height = row_h
# Header
for c, head in enumerate(cols):
    cell = table.cell(0, c)
    cell.text = ""
    cell.fill.solid(); cell.fill.fore_color.rgb = INDIGO
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    r = p.add_run(); r.text = head
    r.font.name = "Calibri"; r.font.size = Pt(13)
    r.font.bold = True; r.font.color.rgb = WHITE
# Body
for ri, row in enumerate(rows, start=1):
    is_lawris = row[0] == "Lawris"
    for ci, val in enumerate(row):
        cell = table.cell(ri, ci)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = INDIGO_D if is_lawris else WHITE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        r = p.add_run(); r.text = val
        r.font.name = "Calibri"; r.font.size = Pt(12)
        r.font.bold = (ci == 0) or is_lawris
        if is_lawris:
            r.font.color.rgb = WHITE
        elif val == "Yes":
            r.font.color.rgb = EMERALD
        elif val == "—":
            r.font.color.rgb = STONE_500
        else:
            r.font.color.rgb = AMBER

add_text(s,
"All five columns checked is a defensible position — not a feature.",
         Inches(0.75), Inches(6.45), Inches(12), Inches(0.6),
         size=18, bold=True, color=INDIGO_D, align=PP_ALIGN.CENTER)

add_notes(s,
"Five columns, not four — the hybrid RAG column is the newest. Harvey does "
"keyword search, IndianKanoon is search only without structured output. "
"We're the only product that does stratified semantic + keyword-boosted "
"retrieval across Indian statutes with source-verification links.")

# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — Roadmap + Ask
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, STONE_50)
add_footer(s, 15)
add_section_label(s, "Roadmap & ask")
add_title(s, "Working product. Clear wedge. Choosing where to go next.",
          size=26)

buckets = [
    ("NEXT 30 DAYS",
     ["Multi-user Supabase RLS cutover",
      "OCR for scanned PDFs (Tesseract)",
      "DOCX/PDF export of drafts",
      "Vercel production + custom domain",
      "IndianKanoon API verification layer"],
     EMERALD),
    ("NEXT 90 DAYS",
     ["React Native mobile shell",
      "Hindi + Marathi UI",
      "Voice-to-hearing-log (Whisper)",
      "Bench-aware outcome prediction",
      "Team collaboration features"],
     BLUE),
    ("THE MARKET",
     ["1.4M registered advocates in India",
      "₹500 / lawyer / month target",
      "₹84 cr ARR at 1% adoption",
      "Beachhead: solo + 2–10 partner firms",
      "Distribution: Bar Council partnerships"],
     INDIGO),
]
for i, (title, items, color) in enumerate(buckets):
    x = Inches(0.75 + i * 4.15)
    y = Inches(2.2)
    add_rounded_box(s, x, y, Inches(3.95), Inches(3.6), fill=WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                             Inches(3.95), Inches(0.1))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = color
    add_text(s, title, x + Inches(0.2), y + Inches(0.2),
             Inches(3.7), Inches(0.4),
             size=13, bold=True, color=color)
    for j, item in enumerate(items):
        iy = y + Inches(0.75 + j * 0.55)
        add_text(s, f"•  {item}",
                 x + Inches(0.3), iy,
                 Inches(3.55), Inches(0.5),
                 size=12, color=STONE_900)

# Ask block
ask = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         Inches(0.75), Inches(6.1),
                         Inches(12.2), Inches(0.95))
ask.line.fill.background()
ask.fill.solid()
ask.fill.fore_color.rgb = INDIGO_D
tf = ask.text_frame
tf.margin_left = Inches(0.4); tf.margin_right = Inches(0.4)
tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.2)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = ("THE ASK   We're looking for honest feedback on what would make this "
          "indispensable to a working lawyer — and a path to the next 1,000 users.")
r.font.name = "Calibri"; r.font.size = Pt(15); r.font.bold = True
r.font.color.rgb = WHITE

add_notes(s,
"Close: 'We have a working product, a clear wedge, a market that wants this "
"badly. We're looking for honest feedback on what would make this indispensable "
"to a working lawyer — and a path to the next thousand users. Thank you.'\n\n"
"DO NOT say 'Questions?' — judges will ask without prompting.\n"
"Time budget: 4:45 of 5:00. Q&A fills the buffer.")

# ─── Save ─────────────────────────────────────────────────────────
out = "/home/tbuser/vit_hackathon/presentation/Lawris_Full.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
