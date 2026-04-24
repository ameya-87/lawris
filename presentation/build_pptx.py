"""Generate Lawris.pptx — 6-slide Round 1 hackathon pitch deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Brand palette ────────────────────────────────────────────────
INDIGO   = RGBColor(0x3F, 0x3D, 0x99)   # primary
INDIGO_D = RGBColor(0x1E, 0x1B, 0x4B)   # dark variant for body text
STONE_50 = RGBColor(0xFA, 0xFA, 0xF9)
STONE_900 = RGBColor(0x1C, 0x19, 0x17)
STONE_500 = RGBColor(0x78, 0x71, 0x6C)
RED      = RGBColor(0xB9, 0x1C, 0x1C)
AMBER    = RGBColor(0xB4, 0x53, 0x09)
EMERALD  = RGBColor(0x05, 0x96, 0x69)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# ─── Setup 16:9 deck ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ─── Helpers ──────────────────────────────────────────────────────
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    return bg

def add_text(slide, text, left, top, width, height,
             size=20, bold=False, color=STONE_900, align=PP_ALIGN.LEFT,
             font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf

def add_pill(slide, text, left, top, fill, fg=WHITE, size=11):
    h = Inches(0.32); w = Inches(1.6)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    pill.adjustments[0] = 0.5
    pill.line.fill.background()
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    tf = pill.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    return pill

def add_footer(slide, n, total=6):
    add_text(slide, "LAWRIS", Inches(0.5), Inches(7.05),
             Inches(2), Inches(0.3), size=10, bold=True, color=STONE_500)
    add_text(slide, f"{n} / {total}", Inches(11.8), Inches(7.05),
             Inches(1), Inches(0.3), size=10, color=STONE_500, align=PP_ALIGN.RIGHT)

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ─── Slide 1 — Title ──────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
add_bg(s1, INDIGO_D)

# big mark
add_text(s1, "lawris.", Inches(0.7), Inches(2.3),
         Inches(8), Inches(2.0),
         size=120, bold=True, color=WHITE, font="Calibri")

# tagline
add_text(s1, "An AI agent for Indian advocates.",
         Inches(0.75), Inches(4.2),
         Inches(11), Inches(0.7),
         size=28, color=RGBColor(0xC7, 0xC4, 0xF0), font="Calibri")

# meta
add_text(s1, "Hackathon — Round 1 Ideation",
         Inches(0.75), Inches(6.5),
         Inches(8), Inches(0.4),
         size=14, color=RGBColor(0x9C, 0x9A, 0xC4))

add_text(s1, "live demo · vit_hackathon",
         Inches(8), Inches(6.5),
         Inches(5), Inches(0.4),
         size=14, color=RGBColor(0x9C, 0x9A, 0xC4),
         align=PP_ALIGN.RIGHT)

add_notes(s1,
"COLD OPEN — say verbatim:\n\n"
"\"An accused person in India has 90 days. If their lawyer misses the chargesheet "
"deadline by even a day, they walk free. Lawyers in this country handle 50 to 200 "
"active cases each — on physical files, WhatsApp, and three different research "
"portals. Calendaring failures are responsible for one in four malpractice claims.\n\n"
"We built an AI agent that makes those misses impossible.\"\n\n"
"Pause. Click to next slide. (Target: 0:30)")

# ─── Slide 2 — The Problem ────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
add_bg(s2, STONE_50)
add_footer(s2, 2)

add_text(s2, "The problem",
         Inches(0.75), Inches(0.6),
         Inches(8), Inches(0.6),
         size=18, bold=True, color=INDIGO)

add_text(s2, "Indian lawyers are drowning in their own workflow.",
         Inches(0.75), Inches(1.1),
         Inches(12), Inches(0.9),
         size=36, bold=True, color=STONE_900)

stats = [
    ("6+ hrs / week", "lost just locating documents",     RED),
    ("40–60%",       "of billable time on repetitive drafts", AMBER),
    ("77%",          "of lawyers run cases out of email", INDIGO),
]
for i, (big, small, color) in enumerate(stats):
    x = Inches(0.75 + i * 4.2)
    add_text(s2, big, x, Inches(2.7), Inches(4), Inches(1.2),
             size=56, bold=True, color=color, font="Calibri")
    add_text(s2, small, x, Inches(3.95), Inches(4), Inches(1.0),
             size=16, color=STONE_500)

add_text(s2,
"No existing Indian tool unifies case management, AI drafting, deadline intelligence, "
"and research grounded in Indian law — MyKase has no AI, Harvey has no Indian law, "
"IndianKanoon is search-only.",
         Inches(0.75), Inches(5.6),
         Inches(12), Inches(1.2),
         size=15, color=STONE_500)

add_notes(s2,
"Read the three stats in 5s each. Don't editorialize each one.\n\n"
"After the third stat, the line: \"No existing tool unifies case management, AI "
"drafting, deadline intelligence, and research grounded in Indian law. We're "
"building it.\"\n\n"
"Click. (Target: 0:30 → 1:00)")

# ─── Slide 3 — The Solution ───────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
add_bg(s3, STONE_50)
add_footer(s3, 3)

add_text(s3, "What we built",
         Inches(0.75), Inches(0.6),
         Inches(8), Inches(0.6),
         size=18, bold=True, color=INDIGO)

add_text(s3, "Lawris",
         Inches(0.75), Inches(1.15),
         Inches(8), Inches(1.1),
         size=56, bold=True, color=STONE_900, font="Calibri")

add_text(s3,
"Open a case → it auto-computes statutory deadlines, drafts court-ready documents "
"from case facts in under 30 seconds, and answers legal questions grounded in the "
"case context.",
         Inches(0.75), Inches(2.5),
         Inches(11.8), Inches(1.5),
         size=22, color=STONE_900)

# Three pillars row
pillars = [
    ("DEADLINE BRAIN",   "60/90-day chargesheet rule, limitation periods, hearing alerts — auto-computed from FIR data."),
    ("AI DRAFTING",      "Bail applications, plaints, written statements — streamed in <30s with real Indian-court formatting."),
    ("CASE-AWARE SEARCH","Research grounded in this case's facts. Citations to Supreme Court precedents and BNSS sections."),
]
for i, (title, body) in enumerate(pillars):
    x = Inches(0.75 + i * 4.2)
    # accent bar
    bar = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(4.5),
                              Inches(0.15), Inches(0.4))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = INDIGO
    add_text(s3, title, Inches(0.75 + i * 4.2 + 0.3), Inches(4.5),
             Inches(3.7), Inches(0.4),
             size=12, bold=True, color=INDIGO)
    add_text(s3, body, Inches(0.75 + i * 4.2 + 0.3), Inches(5.0),
             Inches(3.7), Inches(2.0),
             size=14, color=STONE_900)

add_notes(s3,
"Read the one-line solution. Then: \"Let me show you.\"\n\n"
"Switch immediately to the live browser. Don't dwell on this slide — judges came "
"for the demo, not the bullet points.\n\n"
"Time so far: 1:15. Demo runs 1:15 → 3:30 (135 seconds).")

# ─── Slide 4 — Live Demo placeholder ──────────────────────────────
s4 = prs.slides.add_slide(BLANK)
add_bg(s4, INDIGO_D)
add_footer(s4, 4)

add_text(s4, "Live demo",
         Inches(0.75), Inches(3.1),
         Inches(12), Inches(1.2),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
add_text(s4, "lawris.app  ·  open the POCSO case",
         Inches(0.75), Inches(4.4),
         Inches(12), Inches(0.6),
         size=24, color=RGBColor(0xC7, 0xC4, 0xF0), align=PP_ALIGN.CENTER)

add_notes(s4,
"DEMO SCRIPT (target 2:15 total):\n\n"
"0:00 (1:15) — Dashboard. \"Three matters. ONE CRITICAL: chargesheet on a "
"POCSO case is due in two days. The system computed that from the FIR date and "
"BNSS s.187(3). If she misses it, the accused walks.\"\n\n"
"0:25 (1:40) — Click POCSO case. \"Real FIR, real sections, real client. Watch this.\"\n\n"
"0:35 (1:50) — Documents tab → Draft Bail Application. \"For THIS specific case.\"\n\n"
"0:45–1:05 (2:00–2:20) — STAY QUIET. Let them watch the document write itself.\n\n"
"1:05 (2:20) — \"Court-ready. Cites Satender Kumar Antil, Article 21, FIR number, "
"client name — all from the database. Even spotted the chargesheet was overdue. "
"Twenty-five seconds. No lawyer typed any of this.\"\n\n"
"1:25 (2:40) — Hearings tab → Add Hearing. Type briefly.\n\n"
"1:50 (3:05) — Header AI summary updates. \"Anyone on the team can now read 30s "
"and know where this case stands. Today that lives in one lawyer's head.\"\n\n"
"2:00 (3:15) — Research tab → \"what is default bail under section 187(3) BNSS?\"\n\n"
"2:15 (3:30) — Transition to slide 5.\n\n"
"IF SOMETHING BREAKS: \"Here's the recording from this morning.\" Play backup video. "
"Recover at slide 5.")

# ─── Slide 5 — Why us, not them ───────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
add_bg(s5, STONE_50)
add_footer(s5, 5)

add_text(s5, "Why us, not them",
         Inches(0.75), Inches(0.6),
         Inches(8), Inches(0.6),
         size=18, bold=True, color=INDIGO)

add_text(s5, "Adjacent products solve part of this. None solve all of it.",
         Inches(0.75), Inches(1.15),
         Inches(12), Inches(0.7),
         size=24, bold=True, color=STONE_900)

# Comparison table
cols = ["", "Case mgmt", "AI drafting", "Indian law", "Deadline brain"]
rows = [
    ("Clio (US)",       "Yes",  "Basic", "—",   "Yes"),
    ("MyKase (India)",  "Yes",  "—",     "Yes", "Yes"),
    ("Harvey AI (US)",  "—",    "Yes",   "—",   "—"),
    ("IndianKanoon",    "—",    "—",     "Yes", "—"),
    ("Lawris",          "Yes",  "Yes",   "Yes", "Yes"),
]

table_left = Inches(0.75)
table_top  = Inches(2.4)
col_widths = [Inches(3.1), Inches(2.3), Inches(2.3), Inches(2.3), Inches(2.3)]
row_h = Inches(0.55)

table = s5.shapes.add_table(len(rows) + 1, len(cols),
                            table_left, table_top,
                            sum(col_widths, Emu(0)),
                            row_h * (len(rows) + 1)).table
for i, w in enumerate(col_widths):
    table.columns[i].width = w
for r in range(len(rows) + 1):
    table.rows[r].height = row_h

# Header
for c, head in enumerate(cols):
    cell = table.cell(0, c)
    cell.text = ""
    cell.fill.solid(); cell.fill.fore_color.rgb = INDIGO
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    r = p.add_run(); r.text = head
    r.font.name = "Calibri"; r.font.size = Pt(13)
    r.font.bold = True; r.font.color.rgb = WHITE

# Body rows
for ri, row in enumerate(rows, start=1):
    is_lawris = row[0] == "Lawris"
    for ci, val in enumerate(row):
        cell = table.cell(ri, ci)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = INDIGO_D if is_lawris else WHITE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        r = p.add_run(); r.text = val
        r.font.name = "Calibri"; r.font.size = Pt(13)
        r.font.bold = (ci == 0) or (val == "Yes" and is_lawris) or is_lawris
        if is_lawris:
            r.font.color.rgb = WHITE
        elif val == "Yes":
            r.font.color.rgb = EMERALD
        elif val == "—":
            r.font.color.rgb = STONE_500
        else:
            r.font.color.rgb = AMBER

add_text(s5,
"All four columns checked is a defensible position — not a feature.",
         Inches(0.75), Inches(6.4),
         Inches(12), Inches(0.6),
         size=18, bold=True, color=INDIGO_D, align=PP_ALIGN.CENTER)

add_notes(s5,
"\"Every adjacent product solves PART of this. Clio is great in America but doesn't "
"know Indian law. Harvey is brilliant at drafting but doesn't know Indian law either. "
"MyKase manages cases in India but has no AI. We are the first product where all "
"four columns are checked. That's a defensible position, not a feature.\"\n\n"
"Click. (Target: 3:30 → 4:15)")

# ─── Slide 6 — Roadmap + Ask ──────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
add_bg(s6, STONE_50)
add_footer(s6, 6)

add_text(s6, "Roadmap & ask",
         Inches(0.75), Inches(0.6),
         Inches(8), Inches(0.6),
         size=18, bold=True, color=INDIGO)

add_text(s6, "We have a working product. We're choosing where to go next.",
         Inches(0.75), Inches(1.15),
         Inches(12), Inches(0.7),
         size=24, bold=True, color=STONE_900)

# Three-column roadmap
buckets = [
    ("Next 30 days",
     ["IndianKanoon citation grounding",
      "Document upload + OCR",
      "Multi-user with row-level security"]),
    ("Next 90 days",
     ["Mobile (React Native)",
      "Hindi & Marathi UI",
      "Voice-to-hearing-log",
      "Bench-aware outcome prediction"]),
    ("The market",
     ["1.4M registered advocates in India",
      "₹500 / lawyer / month target",
      "₹84 cr ARR at 1% adoption"]),
]
for i, (title, items) in enumerate(buckets):
    x = Inches(0.75 + i * 4.2)
    add_text(s6, title, x, Inches(2.4), Inches(3.7), Inches(0.5),
             size=15, bold=True, color=INDIGO)
    bullets = "\n".join(f"•  {it}" for it in items)
    add_text(s6, bullets, x, Inches(2.95), Inches(3.7), Inches(3.2),
             size=14, color=STONE_900)

# Ask block
ask_box = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.75), Inches(6.05),
                              Inches(11.8), Inches(0.85))
ask_box.line.fill.background()
ask_box.fill.solid()
ask_box.fill.fore_color.rgb = INDIGO_D
ask_tf = ask_box.text_frame
ask_tf.margin_left = Inches(0.4); ask_tf.margin_right = Inches(0.4)
ask_tf.margin_top = Inches(0.15); ask_tf.margin_bottom = Inches(0.15)
ap = ask_tf.paragraphs[0]
ap.alignment = PP_ALIGN.LEFT
ar = ap.add_run()
ar.text = ("THE ASK   We're looking for honest feedback on what would make this "
           "indispensable to a working lawyer.")
ar.font.name = "Calibri"; ar.font.size = Pt(15); ar.font.bold = True
ar.font.color.rgb = WHITE

add_notes(s6,
"Read the three columns at pace. Then read the ask block.\n\n"
"End on: \"We have a working product, a clear wedge, and a market that wants this "
"badly. We're looking for honest feedback on what would make this indispensable "
"to a working lawyer. Thank you.\"\n\n"
"DO NOT say \"Questions?\" — judges will ask without prompting.\n\n"
"Total target: 4:45 of the 5:00 slot. Q&A fills the buffer.")

# ─── Save ─────────────────────────────────────────────────────────
out = "/home/tbuser/vit_hackathon/presentation/Lawris.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
